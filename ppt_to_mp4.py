import os 
# os.chdir('unoconv')


import tempfile 
import argparse 
from subprocess import call,check_call 
import glob 
import tqdm 
from pptx import Presentation 
from gtts import gTTS 
import pyttsx3 
engine = pyttsx3.init('dummy')
from pdf2image import convert_from_path, convert_from_bytes
from moviepy.editor import VideoFileClip, concatenate_videoclips 
# from moviepy import VideoFileClip, concatenate_videoclips 

import ffmpeg 
import traceback 
from PIL import Image
import PyPDF2
import comtypes.client

from pathlib import Path
import win32com.client


def ppt2pdf(ppt_target_file):
    file_path = Path(ppt_target_file).resolve()
    out_file = file_path.parent / file_path.stem
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
    pdf = powerpoint.Presentations.Open(file_path, WithWindow=False)
    pdf.SaveAs(out_file, 32)
    pdf.Close()
    powerpoint.Quit()


FFMPEG_NAME = 'ffmpeg'

def ppt_presenter (pptx_path, pdf_path, output_path):
  with tempfile.TemporaryDirectory() as temp_path:
    temp_path='' 
    images_from_path = convert_from_path(pdf_path, 500, poppler_path = r'E:/ffff/poppler-0.68.0/bin')
    prs = Presentation(pptx_path)
    # assert len(images_from_path) == len(prs.slides)
    len(images_from_path) == len(prs.slides)
    for i, (slide, image) in enumerate(zip(prs.slides, images_from_path)): 
      if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text 
        tts = gTTS(text=notes, lang='en') 
        image_path = os.path.join(temp_path, 'frame_{}.jpg'.format(i))
        audio_path = os.path.join(temp_path, 'frame_{}.mp3'.format(i))
        image = image.resize((2000,1500), Image.ANTIALIAS)
        image.save(image_path) 
        print(image_path)
        tts.save(audio_path) 
        print(audio_path)
        ffmpeg_call(image_path, audio_path, temp_path, i)

    video_list = [os.path.join(temp_path, 'frame_{}.mp4'.format(i)) \
                 for i in range(len(images_from_path))]
    print(video_list) 
    video_list_str = 'concat:' + '|'.join(video_list)
    ffmpeg_concat(video_list, output_path)

def ffmpeg_call(image_path, audio_path, temp_path, i):
  out_path_mp4 = os.path.join(temp_path, 'frame_{}.mp4'.format(i))
  out_path_ts = os.path.join(temp_path, 'frame_{}.ts'.format(i))
  input_still = ffmpeg.input(image_path)
  input_audio = ffmpeg.input(audio_path) 
  try:
    (
      ffmpeg 
      .concat(input_still, input_audio, v=1, a=1)
      .output(out_path_mp4) 
      .run(overwrite_output=True)
    )
  except Exception: 
    print(traceback.format_exc())
    raise ValueError('A very specific bad thing happened.')

def ffmpeg_concat(video_list_str, out_path):
  print('{}'.format(video_list_str)) 
  vids=[] 
  for i in video_list_str:
    vids.append(VideoFileClip(i))
  # Concat them 
    final = concatenate_videoclips(vids)
  # Write output to the file 
  final.write_videofile(out_path)



# uploaded_filename = 'demo_notes.pptx'

# allvideo = os.listdir("ppt_files")
'''allvideo =  os.listdir('E:/ffff/ppt_files')
for files in allvideo:
    def main():
        # convert_to_pdf("./ppt_files") 
        ppt2pdf('./ppt_files/cybercrime.pptx')
        ppt_presenter(files, files.rsplit(".", 1 )[ 0 ]+".pdf", 
        files.rsplit(".", 1 )[ 0 ]+".mp4")
        files.save(files.rsplit(".", 1 )[ 0 ]+".mp4")
        files.download(files.rsplit(".", 1 )[ 0 ]+".mp4")

if __name__ == '__main__':
  main()'''
  